from pathlib import Path
import nicpy.neh_misc as nm
from datetime import datetime
import win32com.client
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Output selection of slides from multiple files to a single new file
def grab_slides(files_slides_messages, file_slide_sizes, output_directory):

    if not Path(output_directory).is_dir(): raise Exception('Must input a directory in which to output the slide files.')
    dt_string = nm.get_YYYYMMDDHHMMSS_string(datetime.now(), '-', '_')
    output_directory = Path(output_directory)/dt_string
    output_directory.mkdir()

    sizes_and_files = {}
    for file_path in files_slides_messages.keys():
        if file_slide_sizes[file_path] in sizes_and_files.keys():
            sizes_and_files[file_slide_sizes[file_path]].append(file_path)
        else:
            sizes_and_files[file_slide_sizes[file_path]] = [file_path]

    ppt_instance_old = win32com.client.Dispatch('PowerPoint.Application')
    ppt_instance_new = win32com.client.Dispatch('PowerPoint.Application')

    for size in sizes_and_files.keys():

        ########### Copy the first file as the new slide destination file ###########
        first_path = sizes_and_files[size][0]
        new_prs = Presentation(first_path)

        # Helps to prevent loading of the same file's slide master multiple times
        design_index_cache = {}

        # Add successful search details messages to slide notes of first presentation
        for index_slide, message in files_slides_messages[first_path]:
            if new_prs.slides[index_slide].has_notes_slide:
                notes_text = new_prs.slides[index_slide].notes_slide.notes_text_frame.text
                new_prs.slides[index_slide].notes_slide.notes_text_frame.text = message + notes_text
            else:
                new_prs.slides[index_slide].notes_slide.notes_text_frame.text = message

        # Delete the slides not wanted
        first_path_indices = [files_slides_messages[first_path][ind][0] for ind in range(len(files_slides_messages[first_path]))]
        delete_slides = [i for i in range(len(new_prs.slides)) if i not in first_path_indices]
        new_prs = delete_slides_by_indices(new_prs, delete_slides)

        file_name = 'w={},h={}_{}.pptx'.format(size[0], size[1], dt_string)
        file_path = str(Path(output_directory)/file_name)
        new_prs.save(file_path)

        # Then open it with a windowless Powerpoint instance
        read_only, has_title, window = False, False, False
        new_prs = ppt_instance_new.Presentations.open(file_path, read_only, has_title, window)

        ########### For each original presentation other than the one already dealt with... ###########
        for path_to_presentation in sizes_and_files[size]:
            if path_to_presentation == first_path:
                continue
            else:
                if not isinstance(files_slides_messages, dict) or not Path(path_to_presentation).exists() or not isinstance(files_slides_messages[path_to_presentation], list):
                    raise Exception('files_slides_messages must be a dictionary with keys as filepaths and values as lists of corresponding slide indices to grab.')
                # ... open it with another windowless Powerpoint instance
                old_prs = ppt_instance_old.Presentations.open(path_to_presentation, read_only, has_title, window)

                # ... and copy-paste each requested slide from old to new (some pastes fail for unknown reason, so repeat up to 10 times)
                for slide_number_zero_indexing, message in files_slides_messages[path_to_presentation]:
                    insert_index = -1
                    repeat = 1
                    while repeat < 10:
                        try:
                            old_prs.Slides(slide_number_zero_indexing + 1).Copy()
                            new_prs.Slides.Paste(Index=insert_index)
                            repeat = 10
                        except:
                            repeat += 1
                    if path_to_presentation not in design_index_cache.keys():
                        new_prs.Designs.Load(path_to_presentation)
                        design_index_cache[path_to_presentation] = len(list(new_prs.Designs))-1
                    new_prs.Slides(len(new_prs.Slides)).Design = list(new_prs.Designs)[design_index_cache[path_to_presentation]]

                    # append search messages to each slide's notes
                    notes_text = new_prs.Slides(len(new_prs.Slides)).NotesPage.Shapes.Placeholders(2).TextFrame.TextRange.Text
                    new_prs.Slides(len(new_prs.Slides)).NotesPage.Shapes.Placeholders(2).TextFrame.TextRange.Text = message+notes_text
                old_prs.Close()

        n_slides = len(new_prs.Slides)
        new_prs.SaveAs(file_path)
        new_prs.Close()
        print('Successfully output {} slides to {}'.format(n_slides, file_path))

    # Kill ppt_instances
    ppt_instance_new.Quit()
    del ppt_instance_new
    # ppt_instance_old.Quit()
    del ppt_instance_old


def delete_slides_by_indices(presentation, indices):
    if not all([isinstance(el, int) for el in indices]):
        raise Exception('Must provide a list of indices to delete which are all positive integers =< len(n_slides)-1')
    indices = sorted(set(indices), reverse=True)
    for i in indices:
        rId = presentation.slides._sldIdLst[i].rId
        presentation.part.drop_rel(rId)
        del presentation.slides._sldIdLst[i]

    return presentation



def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

